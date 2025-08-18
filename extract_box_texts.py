# extract_box_texts.py
# Extraction + Normalization + Registry update + Metadata logging
# Supports: pdf, docx, pptx, xlsx/csv, txt, jasp

import os, re, io, json, csv, zipfile, shutil, hashlib
from collections import defaultdict
from typing import Dict, List, Tuple
from datetime import datetime

import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
from tqdm import tqdm

# ---------------- Config ----------------
INPUT_DIR    = r"Dr.Mishra-materials"
OUTPUT_DIR   = r"texts"
REGISTRY_DIR = r"registry"
META_DIR     = r"metadata"
PER_FILE_MD  = os.path.join(META_DIR, "per_file")

CORR_PATH = os.path.join(REGISTRY_DIR, "corrections.json")
PEND_PATH = os.path.join(REGISTRY_DIR, "pending.json")

os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(REGISTRY_DIR, exist_ok=True)
os.makedirs(META_DIR, exist_ok=True)
os.makedirs(PER_FILE_MD, exist_ok=True)

UNSUPPORTED = {
    # Do not parse these as text
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff",
    ".mp4", ".avi", ".mov", ".mkv", ".webm",
    ".mp3", ".wav", ".m4a", ".flac", ".ogg",
    ".zip", ".rar", ".7z",
}
HAS_OCRMYPDF = shutil.which("ocrmypdf") is not None

# ---------------- Registry IO ----------------
def _load(path, default):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return default

def _save(path, obj):
    tmp = path + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(obj, f, indent=2, ensure_ascii=False)
    os.replace(tmp, path)

def load_registry():
    corr = _load(CORR_PATH, {"version":1, "rules":{}})
    pend = _load(PEND_PATH, {"version":1, "candidates":{}})
    return corr, pend

def save_registry(corr, pend):
    _save(CORR_PATH, corr)
    _save(PEND_PATH, pend)

# ---------------- Helpers ----------------
def sanitize_name(s: str) -> str:
    s = re.sub(r"[^\w\-. ]+", "_", s)
    return s.strip()[:200] or "file"

def _has_ctrl(s: str) -> bool:
    for ch in s:
        o = ord(ch)
        if (0 <= o < 32 and ch not in ("\t", "\n")) or o == 127:
            return True
    return False

def _is_noisy(s: str) -> bool:
    if not s:
        return True
    printable = sum(1 for ch in s if (ch.isprintable() or ch in "\t\n"))
    return (printable / max(1, len(s))) < 0.85 or _has_ctrl(s)

def suggest_join_fix(bad: str) -> str:
    return re.sub(r"(\S)\s+(fi|fl|ff|ffi|ffl)\s+(\S)", r"\1\2\3", bad, flags=re.IGNORECASE)

def detect_broken_ligatures(text: str) -> List[Tuple[str,str,int]]:
    out = []
    pat = re.compile(r"\b([A-Za-z]{1,24})\s+(fi|fl|ff|ffi|ffl)\s+([A-Za-z]{1,24})\b")
    for m in pat.finditer(text):
        bad = m.group(0)
        if _is_noisy(bad):
            continue
        fix = suggest_join_fix(bad)
        if bad != fix and not _is_noisy(fix):
            out.append((bad, fix, m.start()))
    return out

def apply_corrections(text: str, rules: Dict[str, Dict]) -> str:
    items = sorted(rules.items(), key=lambda kv: len(kv[0]), reverse=True)
    for key, spec in items:
        rep = spec.get("replacement")
        if spec.get("status") != "approved" or rep is None:
            continue
        if spec.get("regex"):
            try:
                text = re.sub(key, rep, text)
            except re.error:
                text = text.replace(key, rep)
        else:
            text = text.replace(key, rep)
    return text

def add_pending(pend: Dict, key: str, replacement: str, source_name: str, sample: str):
    if _is_noisy(key) or _is_noisy(replacement) or _is_noisy(sample):
        return  # skip garbage
    cand = pend["candidates"].get(key)
    if not cand:
        pend["candidates"][key] = {
            "status": "pending",
            "replacement": replacement,
            "total_count": 1,
            "sources": {source_name: 1},
            "samples": [sample[:200]]
        }
    else:
        cand["total_count"] += 1
        cand["sources"][source_name] = cand["sources"].get(source_name, 0) + 1
        if len(cand["samples"]) < 5 and sample[:200] not in cand["samples"]:
            cand["samples"].append(sample[:200])

# ---------------- Cleaner ----------------
class Cleaner:
    def clean(self, text: str) -> str:
        import unicodedata
        t = text or ""
        t = unicodedata.normalize("NFKC", t)
        t = t.replace("\r\n", "\n").replace("\r", "\n")
        t = t.replace("\u00A0", " ").replace("\u00AD", "")
        t = re.sub(r"([A-Za-z])-\n([A-Za-z])", r"\1\2", t)  # join hyphenated
        t = re.sub(r"\n{3,}", "\n\n", t)                   # collapse blank lines
        t = "\n".join(line.rstrip() for line in t.split("\n"))
        return t.strip()

CLEANER = Cleaner()

# ---------------- Extractors ----------------
def pdf_extract(path: str):
    md = {"type":"pdf"}
    def needs_ocr(p):
        try:
            with fitz.open(p) as d:
                for pg in d:
                    if pg.get_text("text").strip():
                        return False
            return True
        except Exception:
            return True
    used_ocr = False
    ocr_tool = None
    if needs_ocr(path) and HAS_OCRMYPDF:
        try:
            out_pdf = path[:-4] + ".ocr.pdf"
            os.system(f'ocrmypdf --skip-text --deskew "{path}" "{out_pdf}"')
            path = out_pdf
            used_ocr = True
            ocr_tool = "ocrmypdf"
        except Exception:
            pass

    try:
        doc = fitz.open(path)
        md["page_count"] = len(doc)
        lines_all = []
        for pg in doc:
            blocks = pg.get_text("dict")["blocks"]
            blocks.sort(key=lambda b: (b.get("bbox", [0,0,0,0])[1], b.get("bbox", [0,0,0,0])[0]))
            for b in blocks:
                if "lines" not in b: 
                    continue
                for ln in b["lines"]:
                    buf=[]
                    for sp in ln.get("spans", []):
                        s = sp.get("text") or ""
                        if not s: continue
                        s = s.replace("\u00AD","").replace("\u00A0"," ")
                        if buf and buf[-1] and buf[-1][-1].isalnum() and s and s[0].isalnum():
                            buf.append(" ")
                        buf.append(s)
                    txt = "".join(buf).rstrip()
                    if txt:
                        lines_all.append(txt)
                lines_all.append("")  # paragraph break
        doc.close()
        text = CLEANER.clean("\n".join(lines_all))
        md["extraction"] = "blocks"
        if used_ocr: md["ocr_tool"] = ocr_tool
        return text, md
    except Exception as e:
        md["error"] = f"pdf: {e}"
        return "", md

def docx_extract(path: str):
    md = {"type":"docx", "paras":0, "tables":0}
    try:
        doc = Document(path)
        parts=[]
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text); md["paras"]+=1
        for tb in doc.tables:
            md["tables"]+=1
            for row in tb.rows:
                parts.append("\t".join(c.text for c in row.cells))
        return CLEANER.clean("\n".join(parts)), md
    except Exception as e:
        md["error"] = f"docx: {e}"
        return "", md

def pptx_extract(path: str):
    md = {"type":"pptx", "slides":0}
    try:
        prs = Presentation(path)
        parts=[]
        for s in prs.slides:
            md["slides"]+=1
            for shp in s.shapes:
                if hasattr(shp, "text") and shp.text:
                    parts.append(shp.text)
            parts.append("")
        return CLEANER.clean("\n".join(parts)), md
    except Exception as e:
        md["error"] = f"pptx: {e}"
        return "", md

def xlsx_extract(path: str):
    md = {"type":"xlsx", "sheets":0}
    try:
        xl = pd.ExcelFile(path, engine="openpyxl")
        parts=[]
        for name in xl.sheet_names:
            df = pd.read_excel(xl, sheet_name=name, dtype=str).fillna("")
            if df.empty: 
                continue
            md["sheets"]+=1
            parts.append(f"=== Sheet: {name} ===")
            parts.append(",".join(map(str, df.columns)))
            parts.append(df.to_csv(index=False).strip())
        return CLEANER.clean("\n".join(parts)), md
    except Exception as e:
        md["error"] = f"xlsx: {e}"
        return "", md

def csv_extract(path: str):
    md = {"type":"csv"}
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return CLEANER.clean(f.read()), md
    except Exception as e:
        md["error"] = f"csv: {e}"
        return "", md

def txt_like_extract(path: str):
    md = {"type":"text"}
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return CLEANER.clean(f.read()), md
    except Exception as e:
        md["error"] = f"text: {e}"
        return "", md

# -------- JASP extractor --------
def _read_zip_text(zf: zipfile.ZipFile, name: str) -> str:
    raw = zf.read(name)
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("latin-1", "ignore")

def summarize_jasp_json(obj, max_lines=300):
    lines=[]
    def walk(o):
        if len(lines) >= max_lines:
            return
        if isinstance(o, dict):
            # Print interesting scalar fields on one line
            interesting = ["title","name","test","statistic","t","F","p","df","df1","df2","N","mean","sd","se","ci","value","estimate","effect","method","dependent","independent","model"]
            parts=[]
            for k in interesting:
                if k in o:
                    v=o[k]
                    if isinstance(v, (str,int,float,bool)) and v not in ("", None):
                        parts.append(f"{k}={v}")
            if parts:
                lines.append("  - " + ", ".join(parts))
            # Simple table printer if present
            if "columns" in o and "rows" in o and isinstance(o["columns"], list) and isinstance(o["rows"], list):
                cols=[str(c.get("name") or c.get("title") or c.get("key") or idx) for idx,c in enumerate(o["columns"])]
                lines.append("TABLE: " + " | ".join(cols))
                for r in o["rows"][:10]:
                    if isinstance(r, dict):
                        row=[str(r.get(c,"")) for c in cols]
                    elif isinstance(r, list):
                        row=[str(x) for x in r]
                    else:
                        row=[str(r)]
                    lines.append("  " + " | ".join(row))
            for v in o.values():
                walk(v)
        elif isinstance(o, list):
            for it in o[:50]:
                walk(it)
    walk(obj)
    if not lines:
        s=json.dumps(obj, ensure_ascii=False)[:1500]
        return [s]
    return lines[:max_lines]

def jasp_extract(path: str):
    """
    JASP (.jasp) is a ZIP container. We:
    - list files
    - emit embedded CSV/TSV (header + sample rows)
    - summarize embedded JSON objects (analyses/tables)
    """
    md = {"type":"jasp", "files":[]}
    parts=[]
    try:
        with zipfile.ZipFile(path) as z:
            names = z.namelist()
            md["files"] = names

            # Pull data tables first
            for name in names:
                lower=name.lower()
                if lower.endswith((".csv",".tsv",".txt")):
                    try:
                        txt = _read_zip_text(z, name)
                    except Exception:
                        continue
                    txt = CLEANER.clean(txt)
                    # leave header plus first ~200 lines to avoid bloat
                    lines = txt.splitlines()
                    sample = "\n".join(lines[:201])
                    parts.append(f"=== JASP Embedded Data: {name} ===")
                    parts.append(sample)
                    parts.append("")

            # Summarize JSON analyses
            for name in names:
                if name.lower().endswith(".json"):
                    try:
                        jtxt = _read_zip_text(z, name)
                        jobj = json.loads(jtxt)
                    except Exception:
                        continue
                    parts.append(f"=== JASP JSON: {name} (summary) ===")
                    parts.extend(summarize_jasp_json(jobj))
                    parts.append("")

        if not parts:
            parts.append("[No readable text sections found in .jasp archive]")
        text = CLEANER.clean("\n".join(parts))
        return text, md

    except Exception as e:
        md["error"] = f"jasp: {e}"
        return "", md

# ---------------- Main run ----------------
def main():
    corr, pend = load_registry()

    stats = {
        "total":0, "ok":0, "failed":0, "unsupported":0, 
        "by_ext": defaultdict(int), "by_type": defaultdict(int)
    }

    file_list=[]
    for root, _, files in os.walk(INPUT_DIR):
        for fn in files:
            file_list.append(os.path.join(root, fn))
    file_list.sort(key=lambda p: p.lower())

    pbar = tqdm(file_list, desc="Extracting")
    for path in pbar:
        stats["total"] += 1
        base = os.path.basename(path)
        ext  = os.path.splitext(base)[1].lower()
        stats["by_ext"][ext] += 1

        if ext in UNSUPPORTED:
            stats["unsupported"] += 1
            continue

        # Extract
        if ext == ".pdf":
            text, md = pdf_extract(path)
        elif ext == ".docx":
            text, md = docx_extract(path)
        elif ext == ".pptx":
            text, md = pptx_extract(path)
        elif ext in (".xlsx", ".xls"):
            text, md = xlsx_extract(path)
        elif ext == ".csv":
            text, md = csv_extract(path)
        elif ext == ".jasp":
            text, md = jasp_extract(path)
        else:
            text, md = txt_like_extract(path)

        per_file = {
            "source_path": path,
            "filename": base,
            "extension": ext,
            "meta": md
        }

        if not text.strip():
            stats["failed"] += 1
            per_file["status"] = "empty_or_failed"
            with open(os.path.join(PER_FILE_MD, sanitize_name(base) + ".json"), "w", encoding="utf-8") as f:
                json.dump(per_file, f, indent=2, ensure_ascii=False)
            continue

        # Normalize + corrections
        text = CLEANER.clean(text)
        text = apply_corrections(text, corr.get("rules", {}))

        # Detect new ligature issues (clean only)
        for bad, fix, idx in detect_broken_ligatures(text):
            snippet = text[max(0, idx-40): idx+40]
            add_pending(pend, bad, fix, base, snippet)

        # Save .txt
        out_name = sanitize_name(os.path.splitext(base)[0]) + ".txt"
        out_path = os.path.join(OUTPUT_DIR, out_name)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(text)

        per_file["status"] = "ok"
        per_file["text_path"] = out_path
        per_file["char_count"] = len(text)
        stats["ok"] += 1
        stats["by_type"][per_file["meta"].get("type","?")] += 1

        with open(os.path.join(PER_FILE_MD, sanitize_name(base) + ".json"), "w", encoding="utf-8") as f:
            json.dump(per_file, f, indent=2, ensure_ascii=False)

    save_registry(corr, pend)

    # Print summary
    print("\n" + "="*60)
    print("📊 EXTRACTION SUMMARY")
    print("="*60)
    print(f"Total files: {stats['total']}")
    print(f"Successfully extracted: {stats['ok']}")
    print(f"Failed/Empty: {stats['failed']}")
    print(f"Unsupported formats: {stats['unsupported']}")
    print("\n📁 Files by type:")
    for k,v in sorted(stats["by_ext"].items()):
        print(f"  {'✅' if k not in UNSUPPORTED else '⏭️'} {k or '[noext]'}: {v} files")
    print("\nDetected content types:")
    for k,v in sorted(stats["by_type"].items()):
        print(f"  {k}: {v} files")

    # Save a machine-readable summary too
    summary = {
        "timestamp": datetime.utcnow().isoformat()+"Z",
        "stats": stats
    }
    with open(os.path.join(META_DIR, "extraction_summary.json"), "w", encoding="utf-8") as f:
        json.dump(summary, f, indent=2, ensure_ascii=False)
    print(f"\n💾 Summary saved to: {os.path.join(META_DIR, 'extraction_summary.json')}")
    print(f"💾 Pending suggestions: {os.path.join(REGISTRY_DIR, 'pending.json')}")

if __name__ == "__main__":
    main()
